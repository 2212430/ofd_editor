#!/usr/bin/env python3
"""Convert PDF to PPTX (one slide per page, page rendered as image). Invoked by Java backend."""
import io
import sys


def convert_pdf_to_pptx(pdf_path: str, pptx_path: str, resolution: int = 150) -> None:
    try:
        import fitz  # pymupdf
        from pptx import Presentation
    except ImportError:
        raise RuntimeError(
            "依赖未安装。请执行: pip install -r backend/requirements-pdf2pptx.txt"
        )

    doc = fitz.open(pdf_path)
    if doc.page_count == 0:
        doc.close()
        raise RuntimeError("PDF 无有效页面")

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    zoom = resolution / 72.0
    matrix = fitz.Matrix(zoom, zoom)

    for page_index in range(doc.page_count):
        page = doc.load_page(page_index)
        rect = page.rect
        width_emu = int(rect.width * 914400 / 72)
        height_emu = int(rect.height * 914400 / 72)

        if page_index == 0:
            prs.slide_width = width_emu
            prs.slide_height = height_emu

        pix = page.get_pixmap(matrix=matrix, alpha=False)
        img_bytes = pix.tobytes("png")

        slide = prs.slides.add_slide(blank_layout)
        stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(stream, 0, 0, prs.slide_width, prs.slide_height)

    doc.close()
    prs.save(pptx_path)


def main() -> int:
    if len(sys.argv) not in (3, 4):
        print("Usage: pdf_to_pptx.py <input.pdf> <output.pptx> [dpi]", file=sys.stderr)
        return 2
    inp, out = sys.argv[1], sys.argv[2]
    dpi = int(sys.argv[3]) if len(sys.argv) == 4 else 150
    try:
        convert_pdf_to_pptx(inp, out, dpi)
    except Exception as exc:
        print(str(exc), file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
